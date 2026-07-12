#!/usr/bin/env python3
# -*- coding: utf-8 -*-

# Copyright © 2025-2026 Wenze Wei. All Rights Reserved.
#
# This file is part of Encre.
# The Encre project belongs to the Dunimd Team.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# You may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
#
# DISCLAIMER: Users must comply with applicable AI regulations.
# Non-compliance may result in service termination or legal liability.

from __future__ import annotations

"""PowerPoint (.pptx) tool.

Creates and edits presentations (slides, titles, bullet content) and returns
the written file path or a summary.
"""

import importlib
import json
import os
from typing import Any

from encre.tools.base import build_tool


async def _presentation_execute(**kwargs: Any) -> str:
    """Presentation execute.

    Args:
        kwargs: Description of the kwargs parameter.
    """
    action = kwargs.get("action", "")
    file_path = kwargs.get("file_path", "")
    if not file_path and action not in ("create",):
        return "Error: 'file_path' is required."

    if action == "read":
        return _ppt_read(file_path)
    elif action == "extract_text":
        return _ppt_extract_text(file_path)
    elif action == "info":
        return _ppt_info(file_path)
    elif action == "create":
        title = kwargs.get("title", "Presentation")
        slides = kwargs.get("slides", [])
        return _ppt_create(file_path or f"{title.replace(' ', '_')}.pptx", title, slides)
    elif action == "add_slide":
        slide_type = kwargs.get("slide_type", "blank")
        content = kwargs.get("content", "")
        return _ppt_add_slide(file_path, slide_type, content)
    elif action == "list_slides":
        return _ppt_list_slides(file_path)
    else:
        return f"Error: Unknown action '{action}'."


def _ensure_pptx() -> None:
    """Ensure pptx."""
    if importlib.util.find_spec("pptx") is None:
        raise ImportError("python-pptx is required. Install with: pip install python-pptx")


def _ppt_read(file_path: str) -> str:
    """Ppt read.

    Args:
        file_path: Description of the file_path parameter.
    """
    _ensure_pptx()
    from pptx import Presentation
    if not os.path.isfile(file_path):
        return f"Error: File not found: {file_path}"
    try:
        prs = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides):
            parts.append(f"\n--- Slide {i + 1} ---")
            layout_name = slide.slide_layout.name if slide.slide_layout else "unknown"
            parts.append(f"Layout: {layout_name}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            parts.append(t)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        parts.append(" | ".join(cells))
        return "\n".join(parts) if len(parts) > 1 else "(empty presentation)"
    except Exception as e:
        return f"Error reading presentation: {e}"


def _ppt_extract_text(file_path: str) -> str:
    """Ppt extract text.

    Args:
        file_path: Description of the file_path parameter.
    """
    _ensure_pptx()
    from pptx import Presentation
    if not os.path.isfile(file_path):
        return f"Error: File not found: {file_path}"
    try:
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
        return "\n".join(texts) if texts else "(no text content)"
    except Exception as e:
        return f"Error extracting text: {e}"


def _ppt_info(file_path: str) -> str:
    """Ppt info.

    Args:
        file_path: Description of the file_path parameter.
    """
    _ensure_pptx()
    from pptx import Presentation
    if not os.path.isfile(file_path):
        return f"Error: File not found: {file_path}"
    try:
        prs = Presentation(file_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        core_props = prs.core_properties
        return json.dumps({
            "file": os.path.basename(file_path),
            "size_bytes": os.path.getsize(file_path),
            "slides": len(prs.slides),
            "slide_width_emu": slide_width,
            "slide_height_emu": slide_height,
            "slide_width_inches": round(slide_width / 914400, 1),
            "slide_height_inches": round(slide_height / 914400, 1),
            "author": str(core_props.author or ""),
            "title": str(core_props.title or ""),
            "created": str(core_props.created or ""),
            "modified": str(core_props.modified or ""),
        }, indent=2)
    except Exception as e:
        return f"Error reading presentation info: {e}"


def _ppt_create(file_path: str, title: str, slides: list) -> str:
    """Ppt create.

    Args:
        file_path: Description of the file_path parameter.
        title: Description of the title parameter.
        slides: Description of the slides parameter.
    """
    _ensure_pptx()
    from pptx import Presentation

    try:
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = prs.slides[0].shapes.title
        if title_shape:
            title_shape.text = title

        if isinstance(slides, list):
            for slide_data in slides:
                if isinstance(slide_data, str):
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    if slide.shapes.title:
                        slide.shapes.title.text = slide_data
                elif isinstance(slide_data, dict):
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide_title = slide_data.get("title", "")
                    content = slide_data.get("content", "")
                    if slide.shapes.title:
                        slide.shapes.title.text = slide_title
                    if content and slide.placeholders and len(slide.placeholders) > 1:
                        tf = slide.placeholders[1].text_frame
                        tf.text = content

        prs.save(file_path)
        return f"Presentation created: {os.path.abspath(file_path)} ({len(prs.slides)} slides)"
    except Exception as e:
        return f"Error creating presentation: {e}"


def _ppt_add_slide(file_path: str, slide_type: str, content: str) -> str:
    """Ppt add slide.

    Args:
        file_path: Description of the file_path parameter.
        slide_type: Description of the slide_type parameter.
        content: Description of the content parameter.
    """
    _ensure_pptx()
    from pptx import Presentation
    if not os.path.isfile(file_path):
        return f"Error: File not found: {file_path}"
    try:
        prs = Presentation(file_path)
        layout_map = {
            "blank": prs.slide_layouts[6],
            "title": prs.slide_layouts[0],
            "content": prs.slide_layouts[1],
            "two_content": prs.slide_layouts[3],
            "section_header": prs.slide_layouts[2],
        }
        layout = layout_map.get(slide_type, prs.slide_layouts[6])
        slide = prs.slides.add_slide(layout)

        if content and slide.shapes.title:
            slide.shapes.title.text = content.split("\n")[0] if "\n" in content else content

        if content and "\n" in content:
            body_lines = content.split("\n")[1:]
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    tf = shape.text_frame
                    tf.text = ""
                    for i, line in enumerate(body_lines):
                        line = line.strip()
                        if not line:
                            continue
                        if i == 0:
                            tf.text = line
                        else:
                            p = tf.add_paragraph()
                            p.text = line

        prs.save(file_path)
        return f"Slide added to {os.path.basename(file_path)} (layout: {slide_type})"
    except Exception as e:
        return f"Error adding slide: {e}"


def _ppt_list_slides(file_path: str) -> str:
    """Ppt list slides.

    Args:
        file_path: Description of the file_path parameter.
    """
    _ensure_pptx()
    from pptx import Presentation
    if not os.path.isfile(file_path):
        return f"Error: File not found: {file_path}"
    try:
        prs = Presentation(file_path)
        if not prs.slides:
            return "No slides in presentation."
        out = [f"Total slides: {len(prs.slides)}"]
        for i, slide in enumerate(prs.slides):
            layout = slide.slide_layout.name if slide.slide_layout else "unknown"
            shape_count = len(slide.shapes)
            title_text = ""
            if slide.shapes.title:
                title_text = slide.shapes.title.text[:60]
            out.append(f"  Slide {i+1}: layout={layout}, shapes={shape_count}, title={title_text!r}")
        return "\n".join(out)
    except Exception as e:
        return f"Error listing slides: {e}"


EncrePresentationTool = build_tool(
    name="presentation",
    description="""Read, create, and edit PowerPoint (.pptx) presentations.

Actions:
- read: Read full presentation content (all slides, text, tables)
- extract_text: Extract all text from the presentation
- info: Get presentation metadata (slides count, dimensions, author)
- create: Create a new presentation with optional slides
- add_slide: Add a slide to an existing presentation
- list_slides: List all slides with layout and shape info

Slide types for add_slide: blank, title, content, two_content, section_header

Requires: pip install python-pptx""",
    input_schema={
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": ["read", "extract_text", "info", "create", "add_slide", "list_slides"],
                "description": "Operation to perform",
            },
            "file_path": {
                "type": "string",
                "description": "Path to the .pptx file",
            },
            "title": {
                "type": "string",
                "description": "Presentation title (for create action)",
            },
            "slides": {
                "type": "array",
                "description": "List of slide definitions for create: strings or {title, content} dicts",
            },
            "slide_type": {
                "type": "string",
                "enum": ["blank", "title", "content", "two_content", "section_header"],
                "description": "Slide layout type for add_slide",
            },
            "content": {
                "type": "string",
                "description": "Slide content (first line = title, rest = body)",
            },
        },
        "required": ["action"],
    },
    execute=_presentation_execute,
    intents=["data", "general", "research"],
    is_concurrency_safe=lambda data: data.get("action") in ("read", "extract_text", "info", "list_slides"),
    is_destructive=lambda args: args.get("action", "") in ("create", "add_slide"),
    category="docs",
    semantic_type="generate",
)
